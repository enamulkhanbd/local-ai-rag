from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from typing import List, Optional
import os
import tempfile
import shutil
import subprocess
import csv
import json
import mailbox
import re
import wave
import zipfile
import logging
from pathlib import Path
from email import policy
from email.parser import BytesParser
from dotenv import load_dotenv
import numpy as np

# LangChain components
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_openai import ChatOpenAI
from langchain_anthropic import ChatAnthropic
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import PyPDFLoader, WebBaseLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains import RetrievalQA
from langchain.prompts import PromptTemplate
from langchain_core.documents import Document
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from extract_msg import Message as OutlookMessage
from openpyxl import load_workbook
from pptx import Presentation
from rapidocr_onnxruntime import RapidOCR
import imageio_ffmpeg
import whisper
import xlrd

# --- Load API Keys ---
ROOT_ENV_PATH = Path(__file__).resolve().parents[1] / ".env"
load_dotenv(dotenv_path=ROOT_ENV_PATH, override=True)

def _clean_env(name: str) -> str:
    return (os.getenv(name) or "").strip()

KEYS = {
    "Gemini": _clean_env("GOOGLE_API_KEY"),
    "OpenAI": _clean_env("OPENAI_API_KEY"),
    "Claude": _clean_env("ANTHROPIC_API_KEY"),
    "DeepSeek": _clean_env("DEEPSEEK_API_KEY"),
    "OpenRouter": _clean_env("OPENROUTER_API_KEY")
}
PROVIDER_PRIORITY = ["OpenRouter", "Gemini", "OpenAI", "Claude", "DeepSeek"]

app = FastAPI()
logger = logging.getLogger("uvicorn.error")

# Enable CORS for React frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # In production, replace with your frontend URL
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- Global State ---
vector_db = None
doc_count = 0
ocr_engine = None
whisper_model = None

PDF_EXTENSIONS = {".pdf"}
WORD_OOXML_EXTENSIONS = {".docx", ".docm", ".dotx", ".dotm"}
WORD_TEXTUTIL_EXTENSIONS = {".doc", ".dot", ".rtf"}
SPREADSHEET_OOXML_EXTENSIONS = {".xlsx", ".xlsm", ".xltx", ".xltm"}
SPREADSHEET_BINARY_EXTENSIONS = {".xls", ".xlt"}
PRESENTATION_OOXML_EXTENSIONS = {".pptx", ".pptm", ".potx", ".ppsx", ".ppsm"}
PRESENTATION_BINARY_EXTENSIONS = {".ppt", ".pot", ".pps"}
TEXT_EXTENSIONS = {".txt", ".md"}
DELIMITED_EXTENSIONS = {".csv", ".tsv"}
STRUCTURED_TEXT_EXTENSIONS = {".json", ".xml", ".html", ".htm"}
OPEN_DOCUMENT_EXTENSIONS = {".odt", ".ods", ".odp"}
EMAIL_EXTENSIONS = {".eml", ".msg", ".mbox"}
IMAGE_EXTENSIONS = {
    ".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tif", ".tiff"
}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".aac", ".ogg", ".flac"}
VIDEO_EXTENSIONS = {
    ".mp4", ".mov", ".m4v", ".avi", ".mkv", ".webm"
}
ACCEPTED_UPLOAD_EXTENSIONS = sorted(
    {
        *PDF_EXTENSIONS,
        *WORD_OOXML_EXTENSIONS,
        *WORD_TEXTUTIL_EXTENSIONS,
        *SPREADSHEET_OOXML_EXTENSIONS,
        *SPREADSHEET_BINARY_EXTENSIONS,
        *PRESENTATION_OOXML_EXTENSIONS,
        *PRESENTATION_BINARY_EXTENSIONS,
        *TEXT_EXTENSIONS,
        *DELIMITED_EXTENSIONS,
        *STRUCTURED_TEXT_EXTENSIONS,
        *OPEN_DOCUMENT_EXTENSIONS,
        *EMAIL_EXTENSIONS,
    }
    | IMAGE_EXTENSIONS
    | AUDIO_EXTENSIONS
    | VIDEO_EXTENSIONS
)
SUPPORTED_UPLOAD_TEXT = ", ".join(ACCEPTED_UPLOAD_EXTENSIONS)

# --- Helper Functions ---
def get_configured_providers() -> List[str]:
    return [provider for provider in PROVIDER_PRIORITY if KEYS.get(provider)]


def get_llm(provider: str):
    if provider not in KEYS or not KEYS[provider]:
        return None

    if provider == "Gemini":
        return ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.3)
    elif provider == "OpenAI":
        return ChatOpenAI(model="gpt-4o-mini", temperature=0.3)
    elif provider == "Claude":
        return ChatAnthropic(model="claude-3-5-sonnet-20240620", temperature=0.3)
    elif provider == "DeepSeek":
        return ChatOpenAI(
            model="deepseek-chat",
            openai_api_key=KEYS["DeepSeek"],
            openai_api_base="https://api.deepseek.com/v1",
            temperature=0.3
        )
    elif provider == "OpenRouter":
        return ChatOpenAI(
            model="deepseek/deepseek-chat",
            openai_api_key=KEYS["OpenRouter"],
            openai_api_base="https://openrouter.ai/api/v1",
            temperature=0.3
        )
    return None

def get_embeddings():
    return HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")


def classify_llm_error(provider: str, exc: Exception) -> tuple[int, str]:
    message = str(exc).strip() or f"{provider} request failed"
    lowered = message.lower()
    alternatives = [name for name in get_configured_providers() if name != provider]
    suggestion = f" Try switching to {alternatives[0]}." if alternatives else ""

    if any(token in lowered for token in ["resourceexhausted", "quota", "rate limit", "429"]):
        return 429, f"{provider} quota or rate limit exceeded.{suggestion}"

    if any(token in lowered for token in ["api key", "authentication", "unauthorized", "forbidden", "401", "403"]):
        return 401, f"{provider} authentication failed. Check the API key and account access."

    return 500, f"{provider} request failed: {message}"


def raise_llm_http_error(provider: str, exc: Exception) -> None:
    status_code, detail = classify_llm_error(provider, exc)
    raise HTTPException(status_code=status_code, detail=detail)


def get_provider_attempt_order(requested_provider: str) -> List[str]:
    configured = get_configured_providers()

    if requested_provider not in PROVIDER_PRIORITY:
        raise HTTPException(status_code=400, detail="Invalid provider")

    if requested_provider in configured:
        return [requested_provider] + [provider for provider in configured if provider != requested_provider]

    return configured


def extract_llm_text(response) -> str:
    content = getattr(response, "content", response)

    if isinstance(content, str):
        return content

    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, str):
                parts.append(item)
            elif isinstance(item, dict) and item.get("text"):
                parts.append(str(item["text"]))
            elif hasattr(item, "get") and item.get("text"):
                parts.append(str(item.get("text")))
        return "\n".join(part.strip() for part in parts if part and str(part).strip())

    return str(content)


def run_chat_with_provider(provider: str, message: str, mode: str) -> dict:
    llm = get_llm(provider)
    if not llm:
        raise RuntimeError(f"{provider} is not configured")

    response = None
    source = "General Knowledge"

    if vector_db:
        qa_chain = RetrievalQA.from_chain_type(
            llm=llm,
            retriever=vector_db.as_retriever(search_kwargs={"k": 3}),
            return_source_documents=False,
            chain_type_kwargs={"prompt": RAG_PROMPT}
        )
        rag_response = qa_chain.invoke(message)["result"].strip()

        if NOT_FOUND_FLAG not in rag_response:
            response = rag_response
            source = "Knowledge Base"

    if response is None:
        if mode == "Knowledge Base Only":
            response = "I couldn't find the answer in your knowledge base."
        else:
            general_response = llm.invoke(f"Answer concisely.\n\nQuestion: {message}")
            response = extract_llm_text(general_response).strip()
            source = "General Knowledge"

    return {
        "response": response,
        "source": source,
        "provider_used": provider,
    }


def get_ocr_engine():
    global ocr_engine

    if ocr_engine is None:
        ocr_engine = RapidOCR()

    return ocr_engine


def get_whisper_model():
    global whisper_model

    if whisper_model is None:
        whisper_model = whisper.load_model("base")

    return whisper_model


def save_upload_to_tempfile(upload: UploadFile) -> tuple[str, str]:
    filename = upload.filename or "uploaded-file"
    suffix = Path(filename).suffix.lower()

    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        shutil.copyfileobj(upload.file, tmp)
        temp_path = tmp.name

    return temp_path, filename


def create_document(filename: str, ingest_type: str, text: str) -> List[Document]:
    cleaned = text.strip()
    if not cleaned:
        raise ValueError(f"No extractable text found in {filename}")

    return [
        Document(
            page_content=cleaned,
            metadata={
                "source": filename,
                "ingest_type": ingest_type,
            },
        )
    ]


def extract_plain_text(file_path: str, filename: str) -> List[Document]:
    text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
    return create_document(filename, "text", text)


def extract_delimited_text(file_path: str, filename: str, delimiter: str) -> List[Document]:
    rows = []
    with open(file_path, "r", encoding="utf-8", errors="ignore", newline="") as handle:
        reader = csv.reader(handle, delimiter=delimiter)
        for row in reader:
            values = [str(cell).strip() for cell in row if str(cell).strip()]
            if values:
                rows.append(" | ".join(values))

    return create_document(filename, "spreadsheet", "\n".join(rows))


def extract_json_text(file_path: str, filename: str) -> List[Document]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as handle:
        data = json.load(handle)
    text = json.dumps(data, indent=2, ensure_ascii=False)
    return create_document(filename, "structured-text", text)


def extract_markup_text(file_path: str, filename: str, parser: str) -> List[Document]:
    raw = Path(file_path).read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(raw, parser)
    text = soup.get_text("\n")
    return create_document(filename, "structured-text", text or raw)


def extract_pdf_text(file_path: str) -> List[Document]:
    return PyPDFLoader(file_path).load()


def extract_docx_text(file_path: str, filename: str) -> List[Document]:
    doc = DocxDocument(file_path)
    parts = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return create_document(filename, "word", "\n".join(parts))


def extract_textutil_text(file_path: str, filename: str) -> List[Document]:
    if shutil.which("textutil") is None:
        raise RuntimeError(f"textutil is not available to process {filename}")

    result = subprocess.run(
        ["textutil", "-convert", "txt", "-stdout", file_path],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0:
        stderr = (result.stderr or "").strip()
        raise RuntimeError(f"Failed to process {filename}. {stderr}")

    return create_document(filename, "word", result.stdout)


def extract_spreadsheet_text(file_path: str, filename: str) -> List[Document]:
    workbook = load_workbook(file_path, read_only=True, data_only=True)
    sections = []

    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if values:
                rows.append(" | ".join(values))

        if rows:
            sections.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))

    workbook.close()
    return create_document(filename, "spreadsheet", "\n\n".join(sections))


def extract_xls_text(file_path: str, filename: str) -> List[Document]:
    workbook = xlrd.open_workbook(file_path)
    sections = []

    for sheet in workbook.sheets():
        rows = []
        for row_index in range(sheet.nrows):
            values = [str(value).strip() for value in sheet.row_values(row_index) if str(value).strip()]
            if values:
                rows.append(" | ".join(values))

        if rows:
            sections.append(f"Sheet: {sheet.name}\n" + "\n".join(rows))

    return create_document(filename, "spreadsheet", "\n\n".join(sections))


def extract_presentation_text(file_path: str, filename: str) -> List[Document]:
    presentation = Presentation(file_path)
    slides = []

    for index, slide in enumerate(presentation.slides, start=1):
        parts = [f"Slide {index}"]

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    parts.append(text)

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))

        if len(parts) > 1:
            slides.append("\n".join(parts))

    return create_document(filename, "presentation", "\n\n".join(slides))


def extract_binary_strings(file_path: str, filename: str) -> List[Document]:
    data = Path(file_path).read_bytes()
    ascii_matches = [match.decode("latin1", errors="ignore") for match in re.findall(rb"[\x20-\x7E]{4,}", data)]
    utf16_matches = [match.decode("utf-16le", errors="ignore") for match in re.findall(rb"(?:[\x20-\x7E]\x00){4,}", data)]

    seen = set()
    chunks = []
    for item in ascii_matches + utf16_matches:
        normalized = " ".join(item.split())
        if normalized and normalized not in seen:
            seen.add(normalized)
            chunks.append(normalized)

    return create_document(filename, "binary-document", "\n".join(chunks))


def extract_open_document_text(file_path: str, filename: str) -> List[Document]:
    with zipfile.ZipFile(file_path) as archive:
        content_xml = archive.read("content.xml").decode("utf-8", errors="ignore")
    soup = BeautifulSoup(content_xml, "xml")
    return create_document(filename, "open-document", soup.get_text("\n"))


def extract_email_body(message) -> str:
    plain_parts = []
    html_parts = []

    if message.is_multipart():
        for part in message.walk():
            if part.get_content_maintype() == "multipart":
                continue
            if part.get_content_disposition() == "attachment":
                continue

            payload = part.get_payload(decode=True) or b""
            charset = part.get_content_charset() or "utf-8"
            decoded = payload.decode(charset, errors="ignore")
            content_type = part.get_content_type()

            if content_type == "text/plain":
                plain_parts.append(decoded)
            elif content_type == "text/html":
                html_parts.append(BeautifulSoup(decoded, "html.parser").get_text("\n"))
    else:
        payload = message.get_payload(decode=True) or b""
        charset = message.get_content_charset() or "utf-8"
        decoded = payload.decode(charset, errors="ignore")
        if message.get_content_type() == "text/html":
            html_parts.append(BeautifulSoup(decoded, "html.parser").get_text("\n"))
        else:
            plain_parts.append(decoded)

    return "\n\n".join(part.strip() for part in (plain_parts or html_parts) if part.strip())


def extract_eml_text(file_path: str, filename: str) -> List[Document]:
    with open(file_path, "rb") as handle:
        message = BytesParser(policy=policy.default).parse(handle)

    subject = message.get("subject", "").strip()
    body = extract_email_body(message)
    return create_document(filename, "email", "\n".join(part for part in [subject, body] if part))


def extract_msg_text(file_path: str, filename: str) -> List[Document]:
    message = OutlookMessage(file_path)
    try:
        subject = (message.subject or "").strip()
        body = (message.body or "").strip()
    finally:
        message.close()

    return create_document(filename, "email", "\n".join(part for part in [subject, body] if part))


def extract_mbox_text(file_path: str, filename: str) -> List[Document]:
    box = mailbox.mbox(file_path)
    messages = []

    for index, message in enumerate(box, start=1):
        subject = (message.get("subject") or "").strip()
        body = extract_email_body(message)
        combined = "\n".join(part for part in [f"Message {index}", subject, body] if part)
        if combined.strip():
            messages.append(combined)

    return create_document(filename, "email", "\n\n".join(messages))


def extract_image_text(file_path: str, filename: str) -> List[Document]:
    engine = get_ocr_engine()
    result, _ = engine(file_path)
    lines = []

    if result:
        for item in result:
            if len(item) >= 2 and item[1]:
                lines.append(str(item[1]).strip())

    return create_document(filename, "image", "\n".join(line for line in lines if line))


def transcode_media_to_wav(file_path: str, filename: str) -> str:
    ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
    with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
        wav_path = tmp.name

    command = [
        ffmpeg_path,
        "-y",
        "-i",
        file_path,
        "-vn",
        "-acodec",
        "pcm_s16le",
        "-ar",
        "16000",
        "-ac",
        "1",
        wav_path,
    ]

    result = subprocess.run(command, capture_output=True, text=True)
    if result.returncode != 0:
        if os.path.exists(wav_path):
            os.remove(wav_path)
        stderr = (result.stderr or "").strip()
        raise RuntimeError(f"Failed to process media file {filename}. {stderr}")

    return wav_path


def load_wav_audio(file_path: str) -> np.ndarray:
    with wave.open(file_path, "rb") as wav_file:
        frame_count = wav_file.getnframes()
        channel_count = wav_file.getnchannels()
        sample_width = wav_file.getsampwidth()
        sample_rate = wav_file.getframerate()
        audio_bytes = wav_file.readframes(frame_count)

    if sample_width != 2:
        raise ValueError(f"Unsupported audio sample width in {file_path}")

    audio = np.frombuffer(audio_bytes, dtype=np.int16).astype(np.float32) / 32768.0

    if channel_count > 1:
        audio = audio.reshape(-1, channel_count).mean(axis=1)

    if sample_rate != 16000:
        raise ValueError(f"Expected 16kHz audio for {file_path}, received {sample_rate}Hz")

    return audio


def extract_audio_transcript(file_path: str, filename: str, ingest_type: str) -> List[Document]:
    wav_path = transcode_media_to_wav(file_path, filename)

    try:
        audio = load_wav_audio(wav_path)
        model = get_whisper_model()
        result = model.transcribe(audio, fp16=False)
        transcript = (result.get("text") or "").strip()
    finally:
        if os.path.exists(wav_path):
            os.remove(wav_path)

    return create_document(filename, ingest_type, transcript)


def extract_uploaded_documents(upload: UploadFile) -> List[Document]:
    temp_path, filename = save_upload_to_tempfile(upload)
    extension = Path(filename).suffix.lower()

    try:
        if extension in PDF_EXTENSIONS:
            return extract_pdf_text(temp_path)

        if extension in WORD_OOXML_EXTENSIONS:
            return extract_docx_text(temp_path, filename)

        if extension in WORD_TEXTUTIL_EXTENSIONS:
            return extract_textutil_text(temp_path, filename)

        if extension in SPREADSHEET_OOXML_EXTENSIONS:
            return extract_spreadsheet_text(temp_path, filename)

        if extension in SPREADSHEET_BINARY_EXTENSIONS:
            return extract_xls_text(temp_path, filename)

        if extension in PRESENTATION_OOXML_EXTENSIONS:
            return extract_presentation_text(temp_path, filename)

        if extension in PRESENTATION_BINARY_EXTENSIONS:
            return extract_binary_strings(temp_path, filename)

        if extension in TEXT_EXTENSIONS:
            return extract_plain_text(temp_path, filename)

        if extension == ".csv":
            return extract_delimited_text(temp_path, filename, ",")

        if extension == ".tsv":
            return extract_delimited_text(temp_path, filename, "\t")

        if extension == ".json":
            return extract_json_text(temp_path, filename)

        if extension in {".html", ".htm"}:
            return extract_markup_text(temp_path, filename, "html.parser")

        if extension == ".xml":
            return extract_markup_text(temp_path, filename, "xml")

        if extension in OPEN_DOCUMENT_EXTENSIONS:
            return extract_open_document_text(temp_path, filename)

        if extension == ".eml":
            return extract_eml_text(temp_path, filename)

        if extension == ".msg":
            return extract_msg_text(temp_path, filename)

        if extension == ".mbox":
            return extract_mbox_text(temp_path, filename)

        if extension in IMAGE_EXTENSIONS:
            return extract_image_text(temp_path, filename)

        if extension in AUDIO_EXTENSIONS:
            return extract_audio_transcript(temp_path, filename, "audio")

        if extension in VIDEO_EXTENSIONS:
            return extract_audio_transcript(temp_path, filename, "video")

        raise ValueError(
            f"Unsupported file type for {filename}. Supported extensions: {SUPPORTED_UPLOAD_TEXT}"
        )
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

NOT_FOUND_FLAG = "NOT_IN_CONTEXT"
rag_prompt_template = f"""Use the following pieces of context to answer the question at the end.
If the context does NOT contain enough information to answer the question,
respond with EXACTLY the phrase "{NOT_FOUND_FLAG}" and nothing else.

Context:
{{context}}

Question: {{question}}

Answer:"""

RAG_PROMPT = PromptTemplate(
    template=rag_prompt_template,
    input_variables=["context", "question"]
)

# --- Endpoints ---

@app.get("/")
def read_root():
    return {"message": "AI Backend API is running"}

@app.get("/providers")
def get_available_providers():
    return get_configured_providers()

@app.post("/process-knowledge")
async def process_knowledge(
    files: List[UploadFile] = File(None),
    url: Optional[str] = Form(None)
):
    global vector_db, doc_count
    all_data = []
    skipped_files = []
    processed_files = []
    
    try:
        if files:
            for file in files:
                try:
                    docs = extract_uploaded_documents(file)
                    all_data.extend(docs)
                    processed_files.append(file.filename or "uploaded-file")
                except Exception as exc:
                    skipped_files.append(
                        {
                            "name": file.filename or "uploaded-file",
                            "reason": str(exc),
                        }
                    )
        
        if url:
            try:
                loader = WebBaseLoader(url)
                all_data.extend(loader.load())
                processed_files.append(url)
            except Exception as exc:
                skipped_files.append({"name": url, "reason": str(exc)})

        if not all_data:
            if skipped_files:
                reasons = "; ".join(f"{item['name']}: {item['reason']}" for item in skipped_files)
                raise HTTPException(status_code=400, detail=reasons)
            raise HTTPException(status_code=400, detail="No content to process")

        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        chunks = text_splitter.split_documents(all_data)
        
        embeddings = get_embeddings()
        vector_db = FAISS.from_documents(chunks, embeddings)
        doc_count = len(chunks)
        
        return {
            "status": "success",
            "chunks_processed": doc_count,
            "processed_files": processed_files,
            "skipped_files": skipped_files,
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/chat")
async def chat(
    message: str = Form(...),
    provider: str = Form(...),
    mode: str = Form("Smart Hybrid (Docs + AI)")
):
    try:
        provider_attempts = get_provider_attempt_order(provider)
        if not provider_attempts:
            raise HTTPException(status_code=503, detail="No configured providers are available.")

        failures = []

        for attempt_provider in provider_attempts:
            try:
                result = run_chat_with_provider(attempt_provider, message, mode)
                if attempt_provider != provider:
                    result["fallback_from"] = provider
                return result
            except Exception as exc:
                status_code, detail = classify_llm_error(attempt_provider, exc)
                logger.warning("Chat request failed for provider %s: %s", attempt_provider, detail)
                failures.append((status_code, detail))

        final_status = failures[0][0] if failures else 503
        final_detail = "All configured providers failed. " + " ".join(detail for _, detail in failures)
        raise HTTPException(status_code=final_status, detail=final_detail)
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Chat request failed for provider %s", provider)
        raise_llm_http_error(provider, e)

@app.post("/clear")
async def clear_knowledge():
    global vector_db, doc_count
    vector_db = None
    doc_count = 0
    return {"status": "cleared"}

if __name__ == "__main__":
    import uvicorn
    import os
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
